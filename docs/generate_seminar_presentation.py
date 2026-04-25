from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_PATH = Path(r"d:/UIT/NAM3/DOAN1/PROJECT/carevia-platform/docs/seminar_ecommerce_fuzzy_topsis.pptx")

TITLE_COLOR = RGBColor(14, 54, 83)
ACCENT_COLOR = RGBColor(18, 122, 118)
ACCENT_SOFT = RGBColor(228, 244, 242)
HIGHLIGHT_COLOR = RGBColor(225, 109, 54)
TEXT_COLOR = RGBColor(45, 55, 72)
MUTED_COLOR = RGBColor(102, 112, 133)
BACKGROUND_COLOR = RGBColor(248, 250, 252)


def resolve_output_path():
    try:
        OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
        with OUTPUT_PATH.open("ab"):
            pass
        return OUTPUT_PATH
    except PermissionError:
        return OUTPUT_PATH.with_stem(f"{OUTPUT_PATH.stem}_updated")


def set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BACKGROUND_COLOR

    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(0.55),
    )
    band.fill.solid()
    band.fill.fore_color.rgb = TITLE_COLOR
    band.line.fill.background()


def add_footer(slide, text="Nguon: Kumar et al., Results in Control and Optimization (2025)"):
    box = slide.shapes.add_textbox(Inches(0.45), Inches(7.0), Inches(12.2), Inches(0.25))
    frame = box.text_frame
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.size = Pt(10)
    paragraph.font.color.rgb = MUTED_COLOR
    paragraph.alignment = PP_ALIGN.RIGHT


def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.8), Inches(12.0), Inches(0.9))
    frame = title_box.text_frame
    paragraph = frame.paragraphs[0]
    paragraph.text = title
    paragraph.font.size = Pt(26)
    paragraph.font.bold = True
    paragraph.font.color.rgb = TITLE_COLOR

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.55), Inches(11.8), Inches(0.5))
        sub_frame = subtitle_box.text_frame
        sub_paragraph = sub_frame.paragraphs[0]
        sub_paragraph.text = subtitle
        sub_paragraph.font.size = Pt(14)
        sub_paragraph.font.color.rgb = MUTED_COLOR


def add_bullet_slide(prs, title, bullets, subtitle=None, callout=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, title, subtitle)

    body = slide.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(7.7), Inches(4.4))
    frame = body.text_frame
    frame.word_wrap = True
    frame.margin_left = Pt(8)

    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = TEXT_COLOR
        paragraph.space_after = Pt(10)

    if callout:
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(9.0),
            Inches(2.1),
            Inches(3.5),
            Inches(2.2),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT_SOFT
        shape.line.color.rgb = ACCENT_COLOR

        frame = shape.text_frame
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        frame.word_wrap = True
        p = frame.paragraphs[0]
        p.text = callout
        p.font.size = Pt(19)
        p.font.bold = True
        p.font.color.rgb = ACCENT_COLOR
        p.alignment = PP_ALIGN.CENTER

    add_footer(slide)
    return slide


def add_two_panel_slide(prs, title, subtitle, left_heading, left_items, right_heading, right_items, bottom_note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, title, subtitle)

    left = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.8),
        Inches(1.95),
        Inches(5.8),
        Inches(4.6 if bottom_note else 4.95),
    )
    left.fill.solid()
    left.fill.fore_color.rgb = RGBColor(240, 248, 250)
    left.line.color.rgb = ACCENT_COLOR

    right = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(6.75),
        Inches(1.95),
        Inches(5.8),
        Inches(4.6 if bottom_note else 4.95),
    )
    right.fill.solid()
    right.fill.fore_color.rgb = RGBColor(255, 248, 241)
    right.line.color.rgb = HIGHLIGHT_COLOR

    for shape, heading, items in [
        (left, left_heading, left_items),
        (right, right_heading, right_items),
    ]:
        frame = shape.text_frame
        frame.word_wrap = True
        p = frame.paragraphs[0]
        p.text = heading
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        p.space_after = Pt(8)
        for item in items:
            paragraph = frame.add_paragraph()
            paragraph.text = item
            paragraph.font.size = Pt(15)
            paragraph.font.color.rgb = TEXT_COLOR
            paragraph.space_after = Pt(6)

    if bottom_note:
        note = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(0.95),
            Inches(6.05),
            Inches(11.75),
            Inches(0.55),
        )
        note.fill.solid()
        note.fill.fore_color.rgb = ACCENT_SOFT
        note.line.color.rgb = ACCENT_COLOR
        note_frame = note.text_frame
        note_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        note_frame.word_wrap = True
        p = note_frame.paragraphs[0]
        p.text = bottom_note
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        p.alignment = PP_ALIGN.CENTER

    add_footer(slide)
    return slide


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    hero = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.65),
        Inches(1.1),
        Inches(12.0),
        Inches(4.7),
    )
    hero.fill.solid()
    hero.fill.fore_color.rgb = RGBColor(236, 245, 247)
    hero.line.color.rgb = ACCENT_COLOR

    title_frame = hero.text_frame
    title_frame.clear()

    p1 = title_frame.paragraphs[0]
    p1.text = "Evaluating consumers benefits in electronic-commerce using fuzzy TOPSIS"
    p1.font.size = Pt(24)
    p1.font.bold = True
    p1.font.color.rgb = TITLE_COLOR
    p1.alignment = PP_ALIGN.CENTER
    p1.space_after = Pt(16)

    p2 = title_frame.add_paragraph()
    p2.text = "Seminar bai bao khoa hoc"
    p2.font.size = Pt(20)
    p2.font.color.rgb = ACCENT_COLOR
    p2.alignment = PP_ALIGN.CENTER
    p2.space_after = Pt(12)

    p3 = title_frame.add_paragraph()
    p3.text = "Tac gia: Rakesh Kumar va cong su | Tap chi: Results in Control and Optimization | 2025"
    p3.font.size = Pt(14)
    p3.font.color.rgb = MUTED_COLOR
    p3.alignment = PP_ALIGN.CENTER
    p3.space_after = Pt(20)

    p4 = title_frame.add_paragraph()
    p4.text = "Nguoi trinh bay: ........................................"
    p4.font.size = Pt(18)
    p4.font.color.rgb = TEXT_COLOR
    p4.alignment = PP_ALIGN.CENTER
    p4.space_after = Pt(8)

    p5 = title_frame.add_paragraph()
    p5.text = "Lop: ........................................"
    p5.font.size = Pt(18)
    p5.font.color.rgb = TEXT_COLOR
    p5.alignment = PP_ALIGN.CENTER

    tag = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.8),
        Inches(6.05),
        Inches(4.1),
        Inches(0.55),
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = HIGHLIGHT_COLOR
    tag.line.fill.background()
    tag.text_frame.text = "Thong diep chinh: usability la yeu to duoc uu tien cao nhat"
    tag.text_frame.paragraphs[0].font.size = Pt(13)
    tag.text_frame.paragraphs[0].font.bold = True
    tag.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    add_footer(slide, "Nguon tai lieu: backend/Seminar.pdf")


def add_criteria_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Khung danh gia cua bai bao", "5 nhom tieu chi va 6 loi ich duoc xep hang")

    left = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.75),
        Inches(1.95),
        Inches(4.9),
        Inches(4.55),
    )
    left.fill.solid()
    left.fill.fore_color.rgb = RGBColor(240, 248, 250)
    left.line.color.rgb = ACCENT_COLOR

    left_frame = left.text_frame
    left_frame.word_wrap = True
    p = left_frame.paragraphs[0]
    p.text = "5 nhom tieu chi"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    for item in [
        "Usability",
        "Service quality",
        "Information quality",
        "Online trust",
        "Technology acceptance",
    ]:
        bullet = left_frame.add_paragraph()
        bullet.text = item
        bullet.font.size = Pt(18)
        bullet.font.color.rgb = TEXT_COLOR
        bullet.level = 0

    right = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(6.0),
        Inches(1.95),
        Inches(6.2),
        Inches(4.55),
    )
    right.fill.solid()
    right.fill.fore_color.rgb = RGBColor(255, 248, 241)
    right.line.color.rgb = HIGHLIGHT_COLOR

    right_frame = right.text_frame
    right_frame.word_wrap = True
    p = right_frame.paragraphs[0]
    p.text = "6 loi ich duoc danh gia"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    for item in [
        "B1: Chuc nang, hieu qua, tim kiem",
        "B2: Phan hoi dich vu giua khach hang va doanh nghiep",
        "B3: Thong tin huu ich, chinh xac, kip thoi",
        "B4: Co che tuong tac va phan hoi hai chieu",
        "B5: Giao dung san pham nhu mo ta",
        "B6: Cong nghe de tiep can, tiet kiem chi phi",
    ]:
        bullet = right_frame.add_paragraph()
        bullet.text = item
        bullet.font.size = Pt(17)
        bullet.font.color.rgb = TEXT_COLOR
        bullet.level = 0

    add_footer(slide)


def add_method_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Thiet ke nghien cuu va Fuzzy TOPSIS")

    steps = [
        ("Du lieu", "Phong van 21 chuyen gia trong 3 nhom ra quyet dinh"),
        ("Trong so", "Gan muc do quan trong bang thang ngon ngu mo"),
        ("Ma tran mo", "Lap ma tran quyet dinh va chuan hoa du lieu"),
        ("So sanh", "Xac dinh FPIS va FNIS"),
        ("Xep hang", "Tinh closeness coefficient de thu duoc thu tu uu tien"),
    ]

    x = 0.85
    for title, text in steps:
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(2.05),
            Inches(2.3),
            Inches(3.5),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = ACCENT_SOFT if title != "Xep hang" else RGBColor(255, 242, 235)
        card.line.color.rgb = ACCENT_COLOR if title != "Xep hang" else HIGHLIGHT_COLOR

        frame = card.text_frame
        frame.word_wrap = True
        p1 = frame.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(18)
        p1.font.bold = True
        p1.font.color.rgb = TITLE_COLOR
        p1.alignment = PP_ALIGN.CENTER
        p1.space_after = Pt(12)

        p2 = frame.add_paragraph()
        p2.text = text
        p2.font.size = Pt(15)
        p2.font.color.rgb = TEXT_COLOR
        p2.alignment = PP_ALIGN.CENTER

        x += 2.45

    add_footer(slide)


def add_algorithm_detail_slide(prs):
    return add_two_panel_slide(
        prs,
        "Dau vao va ky hieu cua mo hinh",
        "Can noi ro C la tieu chi, B la loi ich duoc xep hang",
        "Du lieu dau vao",
        [
            "21 chuyen gia duoc chia thanh 3 nhom DM1, DM2, DM3",
            "5 tieu chi C1 den C5",
            "6 loi ich B1 den B6",
            "Table 4 ghi cac danh gia ngon ngu cho tung cap Cj - Bi",
        ],
        "Y nghia ky hieu",
        [
            "C1: Usability",
            "C2: Service Quality",
            "C3: Information Quality",
            "C4: Online Trust",
            "C5: Technology Acceptance",
            "B1 den B6 la 6 loi ich duoc dua vao xep hang cuoi cung",
        ],
        "Hay nhan manh: bai bao xep hang B1 den B6, khong phai xep hang C1 den C5.",
    )


def add_step1_weight_slide(prs):
    return add_two_panel_slide(
        prs,
        "Step 1: determined the weightage of criteria",
        "Tac gia bien danh gia ngon ngu thanh trong so mo tam giac",
        "Muc tieu cua Step 1",
        [
            "Gan trong so cho tung tieu chi bang ngon ngu danh gia",
            "Khong dung so cung ngay tu dau vi y kien chuyen gia co do mo ho",
            "Table 3 la bang quy doi tu ngon ngu sang fuzzy triangular numbers",
        ],
        "Thang do ngon ngu",
        [
            "Very Low = (0, 0, 0.1)",
            "Low = (0, 0.1, 0.3)",
            "Medium Low = (0.1, 0.3, 0.5)",
            "Medium = (0.3, 0.5, 0.7)",
            "Medium High = (0.5, 0.7, 0.9)",
            "High = (0.7, 0.9, 1.0)",
            "Very High = (0.9, 1.0, 1.0)",
        ],
        "Y nghia: moi muc danh gia duoc mo hinh hoa bang bo ba (min, trung tam, max).",
    )


def add_step2_matrix_slide(prs):
    return add_two_panel_slide(
        prs,
        "Step 2: Construct fuzzy decision matrix",
        "Day la buoc tong hop y kien DM1, DM2, DM3 thanh mot ma tran mo",
        "Tac gia lam gi?",
        [
            "Lay danh gia ngon ngu cua 3 nhom chuyen gia tu Table 4",
            "Doi moi danh gia thanh fuzzy number",
            "Tong hop lai cho tung cap Cj - Bi",
            "Ket qua sau tong hop la Decision Matrix trong Table 5",
        ],
        "Cong thuc va cach doc",
        [
            "xij = (1/k) x (xij1 + xij2 + xij3 + ...)",
            "xij la gia tri tong hop cho loi ich i theo tieu chi j",
            "k la so danh gia duoc lay trung binh",
            "Moi o trong Table 5 la mot bo so mo tam giac dai dien cho y kien tong hop",
        ],
        "Co the giai thich ngan gon: Step 2 bien y kien roi rac thanh du lieu tinh toan duoc.",
    )


def add_step3_normalize_slide(prs):
    return add_two_panel_slide(
        prs,
        "Step 3: Normalise decision matrix",
        "Sau khi co decision matrix, tac gia dua cac gia tri ve cung thang do",
        "Cong thuc chuan hoa",
        [
            "rij = (aij/cj+, bij/cj+, cij/cj+)",
            "cj+ la gia tri can tren lon nhat cua cot j",
            "Ket qua cua buoc nay nam trong Table 6",
        ],
        "Y nghia cua Step 3",
        [
            "Loai bo khac biet ve thang do giua cac tieu chi",
            "Giup cac cot co the duoc so sanh truc tiep",
            "Gia tri gan 1 cho thay phuong an tot hon theo tieu chi do",
            "Day la buoc tien xu ly truoc khi tim FPIS va FNIS",
        ],
        "Neu khong chuan hoa, viec so sanh khoang cach den nghiem ly tuong se khong cong bang.",
    )


def add_step4_ideal_slide(prs):
    return add_two_panel_slide(
        prs,
        "Step 4: Weightage matrix, FPIS va FNIS",
        "Day la phan cot loi cua TOPSIS: xac dinh diem tot nhat va xau nhat de tham chieu",
        "Tac gia lam gi?",
        [
            "Tu ma tran chuan hoa, tao Weightage Matrix trong Table 7",
            "Xac dinh FPIS = Fuzzy Positive Ideal Solution",
            "Xac dinh FNIS = Fuzzy Negative Ideal Solution",
            "Day la hai moc ly tuong de do khoang cach cua moi loi ich",
        ],
        "Cach giai thich truoc lop",
        [
            "FPIS la diem ly tuong tot nhat",
            "FNIS la diem ly tuong xau nhat",
            "Phuong an nao gan FPIS va xa FNIS hon se duoc danh gia cao hon",
            "Sau buoc nay, bai bao chuyen sang tinh khoang cach cho tung B1 den B6",
        ],
        "Noi ngan gon: TOPSIS khong hoi ai tot tuyet doi, ma hoi ai gan diem tot nhat va xa diem xau nhat nhat.",
    )


def add_formula_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Step 5: Closeness from FPIS and FNIS", "Gia tri closeness coefficient la diem so cuoi cung de xep hang")

    formula = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.95),
        Inches(2.0),
        Inches(5.2),
        Inches(2.3),
    )
    formula.fill.solid()
    formula.fill.fore_color.rgb = ACCENT_SOFT
    formula.line.color.rgb = ACCENT_COLOR
    frame = formula.text_frame
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.word_wrap = True
    p1 = frame.paragraphs[0]
    p1.text = "CCi = d-i / (d+i + d-i)"
    p1.font.size = Pt(24)
    p1.font.bold = True
    p1.font.color.rgb = TITLE_COLOR
    p1.alignment = PP_ALIGN.CENTER
    p2 = frame.add_paragraph()
    p2.text = "CC cang gan 1 thi loi ich do cang gan nghiem tot nhat"
    p2.font.size = Pt(16)
    p2.font.color.rgb = TEXT_COLOR
    p2.alignment = PP_ALIGN.CENTER

    explain = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(6.55),
        Inches(2.0),
        Inches(5.8),
        Inches(3.9),
    )
    explain.fill.solid()
    explain.fill.fore_color.rgb = RGBColor(255, 248, 241)
    explain.line.color.rgb = HIGHLIGHT_COLOR
    explain_frame = explain.text_frame
    explain_frame.word_wrap = True
    p = explain_frame.paragraphs[0]
    p.text = "Cach doc Table 8"
    p.font.size = Pt(21)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    for item in [
        "d+i: khoang cach tu loi ich i den FPIS",
        "d-i: khoang cach tu loi ich i den FNIS",
        "B1 co d+ = 1.95 va d- = 3.21 nen CC xap xi 0.620",
        "Neu mot loi ich xa FNIS va khong qua xa FPIS, no se co hang cao",
    ]:
        paragraph = explain_frame.add_paragraph()
        paragraph.text = item
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = TEXT_COLOR
        paragraph.space_after = Pt(8)

    add_footer(slide)


def add_results_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Step 6: Order of ranking", "Table 9 sap xep thu tu cuoi cung theo closeness coefficient")

    chart_data = CategoryChartData()
    chart_data.categories = ["B1", "B4", "B5", "B6", "B2", "B3"]
    chart_data.add_series("Closeness", (0.620, 0.583, 0.582, 0.530, 0.510, 0.500))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.75),
        Inches(2.0),
        Inches(7.2),
        Inches(4.5),
        chart_data,
    ).chart

    chart.has_legend = False
    chart.value_axis.maximum_scale = 0.7
    chart.value_axis.minimum_scale = 0.0
    chart.value_axis.has_major_gridlines = True
    chart.category_axis.tick_labels.font.size = Pt(12)
    chart.value_axis.tick_labels.font.size = Pt(11)
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = ACCENT_COLOR

    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(8.35),
        Inches(2.05),
        Inches(4.3),
        Inches(4.45),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(255, 248, 241)
    box.line.color.rgb = HIGHLIGHT_COLOR
    frame = box.text_frame
    frame.word_wrap = True

    p = frame.paragraphs[0]
    p.text = "Cach ket luan tu bang rank"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    entries = [
        "B1 dung dau: chuc nang, hieu qua, tim kiem",
        "B4 dung thu 2: phan hoi hai chieu",
        "B5 dung thu 3: giao dung mo ta",
        "B3 xep cuoi: thong tin dung la quan trong, nhung khong phai loi ich duoc uu tien nhat trong mau nay",
    ]
    for entry in entries:
        paragraph = frame.add_paragraph()
        paragraph.text = entry
        paragraph.font.size = Pt(17)
        paragraph.font.color.rgb = TEXT_COLOR
        paragraph.space_after = Pt(8)

    add_footer(slide)


def add_strengths_limits_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Danh gia bai bao")

    left = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.8),
        Inches(2.0),
        Inches(5.8),
        Inches(4.5),
    )
    left.fill.solid()
    left.fill.fore_color.rgb = RGBColor(240, 248, 250)
    left.line.color.rgb = ACCENT_COLOR

    right = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(6.75),
        Inches(2.0),
        Inches(5.8),
        Inches(4.5),
    )
    right.fill.solid()
    right.fill.fore_color.rgb = RGBColor(255, 248, 241)
    right.line.color.rgb = HIGHLIGHT_COLOR

    for shape, heading, items, color in [
        (
            left,
            "Diem manh",
            [
                "De tai gan voi bai toan thuc te cua thuong mai dien tu",
                "Khung tieu chi ro rang va de ung dung",
                "Fuzzy TOPSIS phu hop voi du lieu danh gia mo ho",
            ],
            ACCENT_COLOR,
        ),
        (
            right,
            "Han che",
            [
                "Mau nho: 21 chuyen gia",
                "Pham vi chu yeu o An Do",
                "Ket qua trinh bay chua that nhat quan giua tieu chi va loi ich",
                "Can kiem chung them bang du lieu nguoi dung thuc te",
            ],
            HIGHLIGHT_COLOR,
        ),
    ]:
        frame = shape.text_frame
        frame.word_wrap = True
        p = frame.paragraphs[0]
        p.text = heading
        p.font.size = Pt(21)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        for item in items:
            paragraph = frame.add_paragraph()
            paragraph.text = item
            paragraph.font.size = Pt(17)
            paragraph.font.color.rgb = TEXT_COLOR
            paragraph.space_after = Pt(8)
        shape.line.color.rgb = color

    add_footer(slide)


def add_closing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    center = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(1.2),
        Inches(1.5),
        Inches(10.9),
        Inches(4.9),
    )
    center.fill.solid()
    center.fill.fore_color.rgb = RGBColor(236, 245, 247)
    center.line.color.rgb = ACCENT_COLOR

    frame = center.text_frame
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.word_wrap = True
    p1 = frame.paragraphs[0]
    p1.text = "Ket luan"
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = TITLE_COLOR
    p1.alignment = PP_ALIGN.CENTER
    p1.space_after = Pt(18)

    p2 = frame.add_paragraph()
    p2.text = "Bai bao cho thay usability la yeu to uu tien hang dau khi danh gia loi ich nguoi tieu dung tren website e-commerce."
    p2.font.size = Pt(20)
    p2.font.color.rgb = TEXT_COLOR
    p2.alignment = PP_ALIGN.CENTER
    p2.space_after = Pt(14)

    p3 = frame.add_paragraph()
    p3.text = "Doanh nghiep can toi uu trai nghiem tim kiem, thao tac va co che phan hoi truoc khi ky vong tang chuyen doi."
    p3.font.size = Pt(20)
    p3.font.color.rgb = TEXT_COLOR
    p3.alignment = PP_ALIGN.CENTER
    p3.space_after = Pt(26)

    p4 = frame.add_paragraph()
    p4.text = "Cam on thay/co va cac ban da lang nghe\nQ and A"
    p4.font.size = Pt(22)
    p4.font.bold = True
    p4.font.color.rgb = HIGHLIGHT_COLOR
    p4.alignment = PP_ALIGN.CENTER

    add_footer(slide)


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)
    add_bullet_slide(
        prs,
        "Boi canh va ly do chon de tai",
        [
            "Thuong mai dien tu phat trien manh trong ky nguyen so",
            "Nguoi dung co rat nhieu lua chon, muc do canh tranh rat cao",
            "Website kem trai nghiem se mat khach hang rat nhanh",
            "Can xac dinh yeu to nao tao loi ich lon nhat cho nguoi tieu dung",
        ],
        callout="Bai toan cot loi: lam gi de khach hang o lai va mua hang?",
    )
    add_bullet_slide(
        prs,
        "Muc tieu nghien cuu",
        [
            "Xac dinh cac yeu to anh huong den loi ich cua nguoi tieu dung tren website e-commerce",
            "Danh gia 6 loi ich cu the duoi 5 nhom tieu chi",
            "Su dung Fuzzy TOPSIS de xep hang muc do uu tien",
        ],
        callout="Cau hoi nghien cuu: yeu to nao can duoc uu tien nhat?",
    )
    add_criteria_slide(prs)
    add_method_slide(prs)
    add_algorithm_detail_slide(prs)
    add_step1_weight_slide(prs)
    add_step2_matrix_slide(prs)
    add_step3_normalize_slide(prs)
    add_step4_ideal_slide(prs)
    add_formula_slide(prs)
    add_results_slide(prs)
    add_bullet_slide(
        prs,
        "Ham y quan tri cho doanh nghiep",
        [
            "Uu tien toi uu tim kiem, dieu huong va toc do thao tac",
            "Tang co che phan hoi va tuong tac voi khach hang",
            "Dam bao giao dung san pham de xay dung niem tin",
            "Dau tu cong nghe de dung, de truy cap va hop ly ve chi phi",
        ],
        callout="Thong diep quan tri: UX va niem tin canh tranh truoc ca gia re.",
    )
    add_strengths_limits_slide(prs)
    add_bullet_slide(
        prs,
        "Huong nghien cuu tiep theo",
        [
            "Mo rong mau nguoi dung thuc te thay vi chi dung chuyen gia",
            "So sanh giua quoc gia, khu vuc va nhom san pham",
            "Bo sung cac bien nhu gia, toc do giao hang, vung mien",
            "Kiem chung lai mo hinh bang du lieu hanh vi mua sam thuc te",
        ],
        callout="Huong mo rong hop ly nhat: ket hop du lieu chuyen gia va du lieu nguoi dung.",
    )
    add_closing_slide(prs)

    output_path = resolve_output_path()
    prs.save(output_path)
    return output_path


if __name__ == "__main__":
    created_path = build_presentation()
    print(f"Created: {created_path}")